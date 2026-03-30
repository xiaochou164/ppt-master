#!/usr/bin/env python3
"""
PPT Master - SVG to PPTX Conversion Tool

Batch converts SVG files in a project into a PowerPoint presentation.
Each SVG file corresponds to one slide, with SVG embedded in native vector format.

Usage:
    python3 scripts/svg_to_pptx.py <project_path>
    python3 scripts/svg_to_pptx.py <project_path> -o output.pptx
    python3 scripts/svg_to_pptx.py <project_path> --use-final

Examples:
    python3 scripts/svg_to_pptx.py examples/ppt169_demo
    python3 scripts/svg_to_pptx.py examples/ppt169_demo -o presentation.pptx
    python3 scripts/svg_to_pptx.py examples/ppt169_demo --use-final

Dependencies:
    pip install python-pptx

Notes:
    - SVG is embedded in native vector format, preserving editability
    - Requires PowerPoint 2016+ to display correctly
"""

import sys
import os
import argparse
import re
import zipfile
import shutil
import tempfile
from pathlib import Path
from typing import Optional, Tuple, List
from xml.etree import ElementTree as ET

# Check if python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("Error: python-pptx library is missing")
    print("Please run: pip install python-pptx")
    sys.exit(1)

# Import project utility modules
sys.path.insert(0, str(Path(__file__).parent))
try:
    from project_utils import get_project_info
    from config import CANVAS_FORMATS
except ImportError:
    CANVAS_FORMATS = {
        'ppt169': {'name': 'PPT 16:9', 'dimensions': '1280×720', 'viewbox': '0 0 1280 720'},
    }

    def get_project_info(path):
        return {'format': 'unknown', 'name': Path(path).name}

# Import animation module
try:
    from pptx_animations import create_transition_xml, TRANSITIONS
    ANIMATIONS_AVAILABLE = True
except ImportError:
    ANIMATIONS_AVAILABLE = False
    TRANSITIONS = {}

# Import native shape converter
try:
    from svg_to_shapes import convert_svg_to_slide_shapes
    NATIVE_SHAPES_AVAILABLE = True
except ImportError:
    NATIVE_SHAPES_AVAILABLE = False

# SVG to PNG library detection (for Office compatibility mode)
# Prefer CairoSVG (better rendering quality), fall back to svglib
PNG_RENDERER = None  # 'cairosvg' | 'svglib' | None

try:
    import cairosvg
    PNG_RENDERER = 'cairosvg'
except (ImportError, OSError):
    try:
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM
        PNG_RENDERER = 'svglib'
    except (ImportError, OSError):
        pass

def get_png_renderer_info() -> tuple:
    """Get PNG renderer information"""
    if PNG_RENDERER == 'cairosvg':
        return ('cairosvg', '(full gradient/filter support)', None)
    elif PNG_RENDERER == 'svglib':
        return ('svglib', '(some gradients may be lost)', 'Install cairosvg for better results: pip install cairosvg')
    else:
        return (None, '(not installed)', 'Install via: pip install cairosvg or pip install svglib reportlab')


# EMU conversion constants
EMU_PER_INCH = 914400
EMU_PER_PIXEL = EMU_PER_INCH / 96

# XML namespaces
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'asvg': 'http://schemas.microsoft.com/office/drawing/2016/SVG/main',
}

# Register namespaces
for prefix, uri in NAMESPACES.items():
    ET.register_namespace(prefix, uri)


def get_slide_dimensions(canvas_format: str, custom_pixels: Optional[Tuple[int, int]] = None) -> Tuple[int, int]:
    """Get slide dimensions (in EMU units)"""
    if custom_pixels:
        width_px, height_px = custom_pixels
    else:
        if canvas_format not in CANVAS_FORMATS:
            canvas_format = 'ppt169'

        dimensions = CANVAS_FORMATS[canvas_format]['dimensions']
        match = re.match(r'(\d+)[×x](\d+)', dimensions)
        if match:
            width_px = int(match.group(1))
            height_px = int(match.group(2))
        else:
            width_px, height_px = 1280, 720

    return int(width_px * EMU_PER_PIXEL), int(height_px * EMU_PER_PIXEL)


def get_pixel_dimensions(canvas_format: str, custom_pixels: Optional[Tuple[int, int]] = None) -> Tuple[int, int]:
    """Get canvas pixel dimensions"""
    if custom_pixels:
        return custom_pixels

    if canvas_format not in CANVAS_FORMATS:
        canvas_format = 'ppt169'

    dimensions = CANVAS_FORMATS[canvas_format]['dimensions']
    match = re.match(r'(\d+)[×x](\d+)', dimensions)
    if match:
        return int(match.group(1)), int(match.group(2))
    return 1280, 720


def get_viewbox_dimensions(svg_path: Path) -> Optional[Tuple[int, int]]:
    """Extract pixel dimensions from SVG viewBox (returns integers)"""
    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            content = f.read(2000)

        match = re.search(r'viewBox="([^"]+)"', content)
        if not match:
            return None

        parts = re.split(r'[\s,]+', match.group(1).strip())
        if len(parts) < 4:
            return None

        width = float(parts[2])
        height = float(parts[3])
        if width <= 0 or height <= 0:
            return None

        return int(round(width)), int(round(height))
    except Exception:
        return None


def detect_format_from_svg(svg_path: Path) -> Optional[str]:
    """Detect canvas format from SVG file's viewBox"""
    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            content = f.read(2000)

        match = re.search(r'viewBox="([^"]+)"', content)
        if match:
            viewbox = match.group(1)
            for fmt_key, fmt_info in CANVAS_FORMATS.items():
                if fmt_info['viewbox'] == viewbox:
                    return fmt_key
    except Exception:
        pass
    return None


def convert_svg_to_png(svg_path: Path, png_path: Path, width: int = None, height: int = None) -> bool:
    """
    Convert SVG to PNG

    Args:
        svg_path: SVG file path
        png_path: Output PNG file path
        width: Output width (pixels)
        height: Output height (pixels)

    Returns:
        Whether the conversion was successful
    """
    if PNG_RENDERER is None:
        return False

    try:
        if PNG_RENDERER == 'cairosvg':
            # Use CairoSVG (better rendering quality)
            cairosvg.svg2png(
                url=str(svg_path),
                write_to=str(png_path),
                output_width=width,
                output_height=height
            )
            return True

        elif PNG_RENDERER == 'svglib':
            # Use svglib (lightweight, but limited gradient support)
            drawing = svg2rlg(str(svg_path))
            if drawing is None:
                print(f"  Warning: Unable to parse SVG ({svg_path.name})")
                return False

            # Render to PNG
            renderPM.drawToFile(
                drawing,
                str(png_path),
                fmt="PNG",
                configPIL={'quality': 95}
            )
            return True

    except Exception as e:
        print(f"  Warning: SVG to PNG conversion failed ({svg_path.name}): {e}")
        return False

    return False


def find_svg_files(project_path: Path, source: str = 'output') -> Tuple[List[Path], str]:
    """
    Find SVG files in the project

    Args:
        project_path: Project directory path
        source: SVG source directory
            - 'output': svg_output (original version)
            - 'final': svg_final (post-processed, recommended)
            - or any subdirectory name

    Returns:
        (list of SVG files, actual directory name used)
    """
    # Predefined directory mapping
    dir_map = {
        'output': 'svg_output',
        'final': 'svg_final',
    }

    # Get directory name (supports predefined aliases or direct directory names)
    dir_name = dir_map.get(source, source)
    svg_dir = project_path / dir_name

    if not svg_dir.exists():
        print(f"  Warning: {dir_name} directory does not exist, trying svg_output")
        dir_name = 'svg_output'
        svg_dir = project_path / dir_name

    if not svg_dir.exists():
        # Search directly in the specified directory
        if project_path.is_dir():
            svg_dir = project_path
            dir_name = project_path.name
        else:
            return [], ''

    return sorted(svg_dir.glob('*.svg')), dir_name


def find_notes_files(project_path: Path, svg_files: List[Path] = None) -> dict:
    """
    Find notes files in the project

    Supports two matching modes (mixed matching supported):
    1. Match by filename (priority): notes/01_cover.md corresponds to 01_cover.svg
    2. Match by index (backward compatible): notes/slide01.md corresponds to the 1st SVG

    Args:
        project_path: Project directory path
        svg_files: SVG file list (for filename matching)

    Returns:
        Dictionary where key is SVG filename (without extension) and value is notes content
    """
    notes_dir = project_path / 'notes'
    notes = {}

    if not notes_dir.exists():
        return notes

    svg_stems_mapping = {}
    svg_index_mapping = {}
    if svg_files:
        for i, svg_path in enumerate(svg_files, 1):
            svg_stems_mapping[svg_path.stem] = i
            svg_index_mapping[i] = svg_path.stem

    # Collect all notes file information
    for notes_file in notes_dir.glob('*.md'):
        try:
            with open(notes_file, 'r', encoding='utf-8') as f:
                content = f.read().strip()
            if not content:
                continue

            stem = notes_file.stem

            # Try to extract index (backward compatible with slide01.md format)
            match = re.search(r'slide[_]?(\d+)', stem)
            if match:
                index = int(match.group(1))
                mapped_stem = svg_index_mapping.get(index)
                if mapped_stem:
                    notes[mapped_stem] = content

            # Match by filename (overrides backward compatible format)
            if stem in svg_stems_mapping:
                notes[stem] = content
        except Exception:
            pass

    return notes


def markdown_to_plain_text(md_content: str) -> str:
    """
    Convert Markdown notes to plain text (for PPTX notes)

    Args:
        md_content: Markdown formatted notes content

    Returns:
        Plain text content
    """
    def strip_inline_bold(text: str) -> str:
        # Remove Markdown bold markers while keeping content
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'__(.+?)__', r'\1', text)
        return text

    lines = []
    for line in md_content.split('\n'):
        # Skip heading lines (starting with #)
        if line.startswith('#'):
            # Extract heading text
            text = re.sub(r'^#+\s*', '', line).strip()
            text = strip_inline_bold(text)
            if text:
                lines.append(text)
                lines.append('')  # Empty line
        # Handle list items (starting with -)
        elif line.strip().startswith('- '):
            item_text = line.strip()[2:]
            item_text = strip_inline_bold(item_text)
            lines.append('• ' + item_text)
        # Regular lines
        elif line.strip():
            text = strip_inline_bold(line.strip())
            lines.append(text)
        else:
            lines.append('')

    # Merge consecutive empty lines
    result = []
    prev_empty = False
    for line in lines:
        if line == '':
            if not prev_empty:
                result.append(line)
            prev_empty = True
        else:
            result.append(line)
            prev_empty = False

    return '\n'.join(result).strip()


def create_notes_slide_xml(slide_num: int, notes_text: str) -> str:
    """
    Create notes slide XML

    Args:
        slide_num: Slide number
        notes_text: Notes text (plain text format)

    Returns:
        Notes slide XML string
    """
    # Escape XML special characters
    notes_text = notes_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

    # Convert line breaks to <a:p> paragraphs
    paragraphs = []
    for para in notes_text.split('\n'):
        if para.strip():
            paragraphs.append(f'''<a:p>
              <a:r>
                <a:rPr lang="zh-CN" dirty="0"/>
                <a:t>{para}</a:t>
              </a:r>
            </a:p>''')
        else:
            paragraphs.append('<a:p><a:endParaRPr lang="zh-CN" dirty="0"/></a:p>')

    paragraphs_xml = '\n            '.join(paragraphs) if paragraphs else '<a:p><a:endParaRPr lang="zh-CN" dirty="0"/></a:p>'

    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
         xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
         xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Slide Image Placeholder 1"/>
          <p:cNvSpPr>
            <a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/>
          </p:cNvSpPr>
          <p:nvPr>
            <p:ph type="sldImg"/>
          </p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
      </p:sp>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="3" name="Notes Placeholder 2"/>
          <p:cNvSpPr>
            <a:spLocks noGrp="1"/>
          </p:cNvSpPr>
          <p:nvPr>
            <p:ph type="body" idx="1"/>
          </p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
        <p:txBody>
          <a:bodyPr/>
          <a:lstStyle/>
          {paragraphs_xml}
        </p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr>
    <a:masterClrMapping/>
  </p:clrMapOvr>
</p:notes>'''


def create_notes_slide_rels_xml(slide_num: int) -> str:
    """
    Create notes slide relationship file XML

    Args:
        slide_num: Slide number

    Returns:
        Relationship file XML string
    """
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{slide_num}.xml"/>
</Relationships>'''


def create_slide_xml_with_svg(
    slide_num: int,
    png_rid: str,
    svg_rid: str,
    width_emu: int,
    height_emu: int,
    transition: Optional[str] = 'fade',
    transition_duration: float = 0.5,
    auto_advance: Optional[float] = None,
    use_compat_mode: bool = True
) -> str:
    """
    Create slide XML containing an SVG image

    Args:
        slide_num: Slide number
        png_rid: PNG fallback image relationship ID
        svg_rid: SVG relationship ID
        width_emu: Width (EMU)
        height_emu: Height (EMU)
        transition: Transition effect name
        transition_duration: Transition duration (seconds)
        auto_advance: Auto-advance interval (seconds)
        use_compat_mode: Whether to use compatibility mode (PNG + SVG dual format)
    """
    # Generate transition effect XML
    transition_xml = ''
    if transition and ANIMATIONS_AVAILABLE:
        transition_xml = '\n' + create_transition_xml(
            effect=transition,
            duration=transition_duration,
            advance_after=auto_advance
        )

    # Compatibility mode: PNG main image + SVG extension (Office official recommendation)
    if use_compat_mode:
        blip_xml = f'''<a:blip r:embed="{png_rid}">
            <a:extLst>
              <a:ext uri="{{96DAC541-7B7A-43D3-8B79-37D633B846F1}}">
                <asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" r:embed="{svg_rid}"/>
              </a:ext>
            </a:extLst>
          </a:blip>'''
    else:
        # Pure SVG mode (only supported by newer Office versions)
        blip_xml = f'<a:blip r:embed="{svg_rid}"/>'

    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:pic>
        <p:nvPicPr>
          <p:cNvPr id="2" name="SVG Image {slide_num}"/>
          <p:cNvPicPr>
            <a:picLocks noChangeAspect="1"/>
          </p:cNvPicPr>
          <p:nvPr/>
        </p:nvPicPr>
        <p:blipFill>
          {blip_xml}
          <a:stretch>
            <a:fillRect/>
          </a:stretch>
        </p:blipFill>
        <p:spPr>
          <a:xfrm>
            <a:off x="0" y="0"/>
            <a:ext cx="{width_emu}" cy="{height_emu}"/>
          </a:xfrm>
          <a:prstGeom prst="rect">
            <a:avLst/>
          </a:prstGeom>
        </p:spPr>
      </p:pic>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr>
    <a:masterClrMapping/>
  </p:clrMapOvr>{transition_xml}
</p:sld>'''


def create_slide_rels_xml(png_rid: str, png_filename: str, svg_rid: str, svg_filename: str, use_compat_mode: bool = True) -> str:
    """
    Create slide relationship file XML

    Args:
        png_rid: PNG image relationship ID
        png_filename: PNG filename
        svg_rid: SVG relationship ID
        svg_filename: SVG filename
        use_compat_mode: Whether to use compatibility mode
    """
    if use_compat_mode:
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="{png_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/{png_filename}"/>
  <Relationship Id="{svg_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/{svg_filename}"/>
</Relationships>'''
    else:
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="{svg_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/{svg_filename}"/>
</Relationships>'''


def create_pptx_with_native_svg(
    svg_files: List[Path],
    output_path: Path,
    canvas_format: Optional[str] = None,
    verbose: bool = True,
    transition: Optional[str] = 'fade',
    transition_duration: float = 0.5,
    auto_advance: Optional[float] = None,
    use_compat_mode: bool = True,
    notes: Optional[dict] = None,
    enable_notes: bool = True,
    use_native_shapes: bool = False
) -> bool:
    """
    Create a PPTX file with native SVG

    Args:
        svg_files: List of SVG files
        output_path: Output path
        canvas_format: Canvas format
        verbose: Whether to output detailed information
        transition: Transition effect (fade/push/wipe/split/strips/cover/random, default: fade)
        transition_duration: Transition duration in seconds (default: 0.4)
        auto_advance: Auto-advance interval (seconds)
        use_compat_mode: Use Office compatibility mode (PNG + SVG dual format, enabled by default)
        notes: Notes dictionary, key is slide number, value is notes content
        enable_notes: Whether to enable notes embedding (enabled by default)
        use_native_shapes: Convert SVG to native DrawingML shapes (directly editable)
    """
    if not svg_files:
        print("Error: No SVG files found")
        return False

    # Native shapes mode takes priority over compat mode
    if use_native_shapes:
        if not NATIVE_SHAPES_AVAILABLE:
            print("Error: svg_to_shapes module not available, cannot use native shapes mode")
            print("  Falling back to SVG embedding mode")
            use_native_shapes = False
        else:
            # Native shapes mode doesn't use compat/PNG
            use_compat_mode = False

    # Check compatibility mode dependencies
    renderer_name, renderer_status, renderer_hint = get_png_renderer_info()
    if not use_native_shapes and use_compat_mode and PNG_RENDERER is None:
        print("Warning: No PNG rendering library installed, cannot use compatibility mode")
        print(f"  {renderer_hint}")
        print("  Will use pure SVG mode (may not display in Office LTSC 2021 and similar versions)")
        use_compat_mode = False

    # Auto-detect canvas format or get dimensions from viewBox
    custom_pixels: Optional[Tuple[int, int]] = None
    if canvas_format is None:
        canvas_format = detect_format_from_svg(svg_files[0])
        if canvas_format and verbose:
            format_name = CANVAS_FORMATS.get(canvas_format, {}).get('name', canvas_format)
            print(f"  Detected canvas format: {format_name}")

    if canvas_format is None:
        custom_pixels = get_viewbox_dimensions(svg_files[0])
        if custom_pixels and verbose:
            print(f"  Using SVG viewBox dimensions: {custom_pixels[0]} x {custom_pixels[1]} px")

    if canvas_format is None and custom_pixels is None:
        canvas_format = 'ppt169'
        if verbose:
            print(f"  Using default format: PPT 16:9")

    width_emu, height_emu = get_slide_dimensions(canvas_format or 'ppt169', custom_pixels)
    pixel_width, pixel_height = get_pixel_dimensions(canvas_format or 'ppt169', custom_pixels)

    if verbose:
        print(f"  Slide dimensions: {pixel_width} x {pixel_height} px")
        print(f"  SVG file count: {len(svg_files)}")
        if use_native_shapes:
            print(f"  Mode: Native DrawingML shapes (directly editable)")
        elif use_compat_mode:
            print(f"  Compatibility mode: Enabled (PNG + SVG dual format)")
            print(f"  PNG renderer: {renderer_name} {renderer_status}")
        else:
            print(f"  Compatibility mode: Disabled (pure SVG)")
        if transition:
            trans_name = TRANSITIONS.get(transition, {}).get('name', transition) if TRANSITIONS else transition
            print(f"  Transition effect: {trans_name}")
        if enable_notes and notes:
            print(f"  Speaker notes: {len(notes)} page(s)")
        elif enable_notes:
            print(f"  Speaker notes: Enabled (no notes files found)")
        else:
            print(f"  Speaker notes: Disabled")
        print()

    # Create temporary directory
    temp_dir = Path(tempfile.mkdtemp())

    try:
        # First create base PPTX with python-pptx
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu

        # Add blank slides as placeholders
        blank_layout = prs.slide_layouts[6]
        for _ in svg_files:
            prs.slides.add_slide(blank_layout)

        # Save base PPTX
        base_pptx = temp_dir / 'base.pptx'
        prs.save(str(base_pptx))

        # Extract PPTX
        extract_dir = temp_dir / 'pptx_content'
        with zipfile.ZipFile(base_pptx, 'r') as zf:
            zf.extractall(extract_dir)

        # Create media directory
        media_dir = extract_dir / 'ppt' / 'media'
        media_dir.mkdir(exist_ok=True)

        # Process each SVG file
        success_count = 0
        any_png_generated = False

        for i, svg_path in enumerate(svg_files, 1):
            slide_num = i

            try:
                # ---- Native shapes mode ----
                if use_native_shapes:
                    slide_xml, media_files_dict, rel_entries = convert_svg_to_slide_shapes(
                        svg_path, slide_num=slide_num, verbose=verbose
                    )

                    # Add transition if specified
                    if transition and ANIMATIONS_AVAILABLE:
                        transition_xml = '\n' + create_transition_xml(
                            effect=transition,
                            duration=transition_duration,
                            advance_after=auto_advance
                        )
                        slide_xml = slide_xml.replace(
                            '</p:sld>',
                            transition_xml + '\n</p:sld>'
                        )

                    # Write slide XML
                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    with open(slide_xml_path, 'w', encoding='utf-8') as f:
                        f.write(slide_xml)

                    # Write media files (from embedded images)
                    for media_name, media_data in media_files_dict.items():
                        with open(media_dir / media_name, 'wb') as f:
                            f.write(media_data)

                    # Build relationships XML
                    rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                    rels_dir.mkdir(exist_ok=True)
                    rels_path = rels_dir / f'slide{slide_num}.xml.rels'

                    extra_rels = ''
                    for rel in rel_entries:
                        extra_rels += f'\n  <Relationship Id="{rel["id"]}" Type="{rel["type"]}" Target="{rel["target"]}"/>'

                    rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>{extra_rels}
</Relationships>'''
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(rels_xml)

                    # Track image formats for Content_Types
                    for media_name in media_files_dict:
                        ext = media_name.rsplit('.', 1)[-1].lower()
                        if ext in ('png', 'jpg', 'jpeg', 'gif', 'webp'):
                            any_png_generated = True  # reuse flag for any image

                # ---- Legacy SVG embedding mode ----
                else:
                    svg_filename = f'image{i}.svg'
                    png_filename = f'image{i}.png'
                    png_rid = 'rId2'
                    svg_rid = 'rId3' if use_compat_mode else 'rId2'

                    # Copy SVG to media directory
                    shutil.copy(svg_path, media_dir / svg_filename)

                    # Compatibility mode: generate PNG fallback image
                    slide_has_png = False
                    if use_compat_mode:
                        png_path = media_dir / png_filename
                        png_success = convert_svg_to_png(
                            svg_path,
                            png_path,
                            width=pixel_width,
                            height=pixel_height
                        )
                        if png_success:
                            slide_has_png = True
                            any_png_generated = True
                        else:
                            # PNG generation failed, fall back to pure SVG
                            if verbose:
                                print(f"  [{i}/{len(svg_files)}] {svg_path.name} - PNG generation failed, using pure SVG")
                            svg_rid = 'rId2'

                    # Update slide XML
                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    slide_xml = create_slide_xml_with_svg(
                        slide_num,
                        png_rid=png_rid,
                        svg_rid=svg_rid,
                        width_emu=width_emu,
                        height_emu=height_emu,
                        transition=transition,
                        transition_duration=transition_duration,
                        auto_advance=auto_advance,
                        use_compat_mode=(use_compat_mode and slide_has_png)
                    )
                    with open(slide_xml_path, 'w', encoding='utf-8') as f:
                        f.write(slide_xml)

                    # Create/update relationship file
                    rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                    rels_dir.mkdir(exist_ok=True)
                    rels_path = rels_dir / f'slide{slide_num}.xml.rels'
                    rels_xml = create_slide_rels_xml(
                        png_rid=png_rid,
                        png_filename=png_filename,
                        svg_rid=svg_rid,
                        svg_filename=svg_filename,
                        use_compat_mode=(use_compat_mode and slide_has_png)
                    )
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(rels_xml)

                # Process notes (shared between native and legacy mode)
                notes_content = ''
                if enable_notes:
                    # Match by filename (new logic) or by index (backward compatible)
                    svg_stem = svg_path.stem
                    notes_content = notes.get(svg_stem, '') if notes else ''
                    if notes_content:
                        notes_text = markdown_to_plain_text(notes_content)
                    else:
                        notes_text = ''  # Empty notes

                    # Create notesSlides directory
                    notes_slides_dir = extract_dir / 'ppt' / 'notesSlides'
                    notes_slides_dir.mkdir(exist_ok=True)

                    # Create notes slide XML
                    notes_xml_path = notes_slides_dir / f'notesSlide{slide_num}.xml'
                    notes_xml = create_notes_slide_xml(slide_num, notes_text)
                    with open(notes_xml_path, 'w', encoding='utf-8') as f:
                        f.write(notes_xml)

                    # Create notes slide relationship file
                    notes_rels_dir = notes_slides_dir / '_rels'
                    notes_rels_dir.mkdir(exist_ok=True)
                    notes_rels_path = notes_rels_dir / f'notesSlide{slide_num}.xml.rels'
                    notes_rels_xml = create_notes_slide_rels_xml(slide_num)
                    with open(notes_rels_path, 'w', encoding='utf-8') as f:
                        f.write(notes_rels_xml)

                    # Update slide.xml.rels to add notes association
                    with open(rels_path, 'r', encoding='utf-8') as f:
                        slide_rels_content = f.read()
                    notes_rel = f'  <Relationship Id="rId10" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" Target="../notesSlides/notesSlide{slide_num}.xml"/>'
                    slide_rels_content = slide_rels_content.replace('</Relationships>', notes_rel + '\n</Relationships>')
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(slide_rels_content)

                if verbose:
                    if use_native_shapes:
                        mode_str = " (Native)"
                    elif use_compat_mode and not use_native_shapes:
                        mode_str = " (PNG+SVG)" if any_png_generated else " (SVG)"
                    else:
                        mode_str = " (SVG)"
                    has_notes = enable_notes and bool(notes_content)
                    notes_str = " +notes" if has_notes else ""
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name}{mode_str}{notes_str}")

                success_count += 1

            except Exception as e:
                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name} - Error: {e}")

        # Update [Content_Types].xml to add media types
        content_types_path = extract_dir / '[Content_Types].xml'
        with open(content_types_path, 'r', encoding='utf-8') as f:
            content_types = f.read()

        types_to_add = []
        if not use_native_shapes:
            # Legacy mode: need SVG extension type
            if 'Extension="svg"' not in content_types:
                types_to_add.append('  <Default Extension="svg" ContentType="image/svg+xml"/>')
        if any_png_generated and 'Extension="png"' not in content_types:
            types_to_add.append('  <Default Extension="png" ContentType="image/png"/>')
        # Add jpg support for native mode embedded images
        if use_native_shapes and 'Extension="jpg"' not in content_types:
            types_to_add.append('  <Default Extension="jpg" ContentType="image/jpeg"/>')
        if use_native_shapes and 'Extension="jpeg"' not in content_types:
            types_to_add.append('  <Default Extension="jpeg" ContentType="image/jpeg"/>')

        if types_to_add:
            content_types = content_types.replace(
                '</Types>',
                '\n'.join(types_to_add) + '\n</Types>'
            )
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        # Add notesSlides content types
        if enable_notes:
            for i in range(1, len(svg_files) + 1):
                override = f'  <Override PartName="/ppt/notesSlides/notesSlide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                if override not in content_types:
                    content_types = content_types.replace('</Types>', override + '\n</Types>')
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        # Repackage PPTX
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(extract_dir)
                    zf.write(file_path, arcname)

        if verbose:
            print()
            print(f"[Done] Saved: {output_path}")
            print(f"  Succeeded: {success_count}, Failed: {len(svg_files) - success_count}")
            if use_compat_mode and any_png_generated:
                print(f"  Mode: Office compatibility mode (supports all Office versions)")
                # If using svglib, provide upgrade hint
                if PNG_RENDERER == 'svglib' and renderer_hint:
                    print(f"  [Tip] {renderer_hint}")

        return success_count == len(svg_files)

    finally:
        # Clean up temporary directory
        shutil.rmtree(temp_dir, ignore_errors=True)


def main():
    # Build transition effect option list
    transition_choices = ['none'] + (list(TRANSITIONS.keys()) if TRANSITIONS else ['fade', 'push', 'wipe', 'split', 'strips', 'cover', 'random'])

    parser = argparse.ArgumentParser(
        description='PPT Master - SVG to PPTX Tool (Office Compatibility Mode)',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=f'''
Examples:
    %(prog)s examples/ppt169_demo -s final    # Default: native + SVG reference (two files)
    %(prog)s examples/ppt169_demo --only native   # Only native shapes version
    %(prog)s examples/ppt169_demo --only legacy   # Only SVG image version
    %(prog)s examples/ppt169_demo -o out.pptx     # Explicit path (SVG ref → out_svg.pptx)

    # Disable transition / change transition effect
    %(prog)s examples/ppt169_demo -t none
    %(prog)s examples/ppt169_demo -t push --transition-duration 1.0

SVG source directory (-s):
    output   - svg_output (original version)
    final    - svg_final (post-processed, recommended)
    <any>    - Specify a subdirectory name directly

Transition effects (-t/--transition):
    {', '.join(transition_choices)}

Compatibility mode (enabled by default):
    - Automatically generates PNG fallback images, SVG embedded as extension
    - Compatible with all Office versions (including Office LTSC 2021)
    - Newer Office still displays SVG (editable), older versions display PNG
    - Optional PNG renderer for Office compatibility mode:
      pip install cairosvg
      or: pip install svglib reportlab
    - Use --no-compat to disable (only Office 2019+ supported)

Speaker notes (enabled by default):
    - Automatically reads Markdown notes files from the notes/ directory
    - Supports two naming conventions:
      1. Match by filename (recommended): 01_cover.md corresponds to 01_cover.svg
      2. Match by index: slide01.md corresponds to the 1st SVG (backward compatible)
    - Use --no-notes to disable
'''
    )

    parser.add_argument('project_path', type=str, help='Project directory path')
    parser.add_argument('-o', '--output', type=str, default=None, help='Output file path')
    parser.add_argument('-s', '--source', type=str, default='output',
                        help='SVG source: output/final or any subdirectory name (recommended: final)')
    parser.add_argument('-f', '--format', type=str, choices=list(CANVAS_FORMATS.keys()), default=None, help='Specify canvas format')
    parser.add_argument('-q', '--quiet', action='store_true', help='Quiet mode')

    # Compatibility mode arguments
    parser.add_argument('--no-compat', action='store_true',
                        help='Disable Office compatibility mode (pure SVG only, requires Office 2019+)')

    # Output mode arguments
    mode_group = parser.add_mutually_exclusive_group()
    mode_group.add_argument('--only', type=str, choices=['native', 'legacy'], default=None,
                            help='Only generate one version: native (editable shapes) or legacy (SVG image)')
    mode_group.add_argument('--native', action='store_true', default=False,
                            help='(Deprecated, now default) Convert SVG to native DrawingML shapes')

    # Transition effect arguments
    parser.add_argument('-t', '--transition', type=str, choices=transition_choices, default='fade',
                        help='Page transition effect (default: fade, use "none" to disable)')
    parser.add_argument('--transition-duration', type=float, default=0.4,
                        help='Transition duration in seconds (default: 0.5)')
    parser.add_argument('--auto-advance', type=float, default=None,
                        help='Auto-advance interval in seconds (default: manual advance)')

    # Notes arguments
    parser.add_argument('--no-notes', action='store_true',
                        help='Disable speaker notes embedding (enabled by default)')

    args = parser.parse_args()

    project_path = Path(args.project_path)
    if not project_path.exists():
        print(f"Error: Path does not exist: {project_path}")
        sys.exit(1)

    try:
        project_info = get_project_info(str(project_path))
        project_name = project_info.get('name', project_path.name)
        detected_format = project_info.get('format')
    except Exception:
        project_name = project_path.name
        detected_format = None

    canvas_format = args.format
    if canvas_format is None and detected_format and detected_format != 'unknown':
        canvas_format = detected_format

    svg_files, source_dir_name = find_svg_files(project_path, args.source)

    if not svg_files:
        print("Error: No SVG files found")
        sys.exit(1)

    # Determine which versions to generate
    # Default: both native + SVG reference.  --only limits to one.
    only_mode = args.only  # None | 'native' | 'legacy'
    gen_native = only_mode in (None, 'native')
    gen_legacy = only_mode in (None, 'legacy')

    # --native flag (deprecated) maps to --only native
    if args.native and only_mode is None:
        gen_legacy = False

    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    if args.output:
        # Explicit output: use as-is for the primary version
        output_base = Path(args.output)
        native_path = output_base
        stem = output_base.stem
        legacy_path = output_base.parent / f"{stem}_svg{output_base.suffix}"
    else:
        native_path = project_path / f"{project_name}_{timestamp}.pptx"
        legacy_path = project_path / f"{project_name}_{timestamp}_svg.pptx"

    native_path.parent.mkdir(parents=True, exist_ok=True)

    verbose = not args.quiet

    # Read notes files
    enable_notes = not args.no_notes
    notes = {}
    if enable_notes:
        notes = find_notes_files(project_path, svg_files)

    # Convert 'none' string to None for downstream logic
    transition = args.transition if args.transition != 'none' else None

    # Shared args for both runs
    shared_kwargs = dict(
        svg_files=svg_files,
        canvas_format=canvas_format,
        verbose=verbose,
        transition=transition,
        transition_duration=args.transition_duration,
        auto_advance=args.auto_advance,
        use_compat_mode=not args.no_compat,
        notes=notes,
        enable_notes=enable_notes,
    )

    success = True

    # --- Native shapes version (primary) ---
    if gen_native:
        if verbose:
            print("PPT Master - SVG to PPTX Tool")
            print("=" * 50)
            print(f"  Project path: {project_path}")
            print(f"  SVG directory: {source_dir_name}")
            print(f"  Output file: {native_path}")
            print()

        ok = create_pptx_with_native_svg(
            output_path=native_path,
            use_native_shapes=True,
            **shared_kwargs,
        )
        success = success and ok

    # --- SVG image reference version ---
    if gen_legacy:
        if verbose:
            if gen_native:
                print()
                print("-" * 50)
            print("PPT Master - SVG to PPTX Tool (SVG Reference)")
            print("=" * 50)
            print(f"  Project path: {project_path}")
            print(f"  SVG directory: {source_dir_name}")
            print(f"  Output file: {legacy_path}")
            print()

        ok = create_pptx_with_native_svg(
            output_path=legacy_path,
            use_native_shapes=False,
            **shared_kwargs,
        )
        success = success and ok

    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
